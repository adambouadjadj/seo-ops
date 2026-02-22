#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
tools/make_template.py
----------------------
Crée un template PPT vierge à partir du reporting de référence :
- Structure et mise en forme intactes
- En-têtes de mois mis à jour : janv.-26 → févr.-27 (14 mois)
- Cellules de données vidées

Usage:
    python tools/make_template.py
"""

import unicodedata
from pptx import Presentation
from pptx.oxml.ns import qn

SOURCE_PPT = "C:/Users/abouadjadj/Downloads/Reporting AB Webperf 14-11-2025.pptx"
OUTPUT_PPT = "tools/webperf_template.pptx"

NEW_MONTHS = [
    "janv.-26", "févr.-26", "mars-26",  "avr.-26",
    "mai-26",   "juin-26",  "juil.-26", "août-26",
    "sept.-26", "oct.-26",  "nov.-26",  "déc.-26",
    "janv.-27", "févr.-27",
]

MONTH_PREFIXES = [
    "janv", "fevr", "mars", "avr", "mai", "juin",
    "juil", "aout", "sept", "oct", "nov", "dec",
    "jan", "feb", "mar", "apr", "jul", "aug", "sep",
]


def _strip_accents(text):
    return ''.join(
        c for c in unicodedata.normalize('NFD', text)
        if unicodedata.category(c) != 'Mn'
    )


def looks_like_month(text):
    if not text:
        return False
    clean = _strip_accents(text.lower()).replace('.', '').replace('-', '').strip()
    for prefix in MONTH_PREFIXES:
        if clean.startswith(prefix):
            return True
    return False


def set_cell_text(cell, new_text):
    """Remplace le texte d'une cellule en préservant la mise en forme."""
    t_elems = cell.text_frame._txBody.findall('.//' + qn('a:t'))
    if not t_elems:
        cell.text = new_text
        return
    t_elems[0].text = new_text
    for t_elem in t_elems[1:]:
        t_elem.text = ''


def process_table(table):
    """Met à jour les en-têtes de mois et vide les cellules de données."""
    row0 = table.rows[0]

    # Trouver la première colonne contenant un mois
    first_month_col = -1
    for i, cell in enumerate(row0.cells):
        if looks_like_month(cell.text.strip()):
            first_month_col = i
            break

    if first_month_col == -1:
        return  # Pas de colonne mois → skip

    # Mettre à jour les en-têtes de mois (row 0)
    for i, cell in enumerate(list(row0.cells)[first_month_col:]):
        new_label = NEW_MONTHS[i] if i < len(NEW_MONTHS) else ""
        set_cell_text(cell, new_label)

    # Vider les données (rows 1+, colonnes de mois uniquement)
    for row in list(table.rows)[1:]:
        for cell in list(row.cells)[first_month_col:]:
            set_cell_text(cell, "")


def main():
    print(f"Source  : {SOURCE_PPT}")
    print(f"Output  : {OUTPUT_PPT}")

    prs = Presentation(SOURCE_PPT)

    table_count = 0
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 19:  # Table
                process_table(shape.table)
                table_count += 1
                print(f"  Slide {slide_idx:2d} — table {len(shape.table.rows)}r x {len(shape.table.columns)}c vidée")

    prs.save(OUTPUT_PPT)
    print(f"\nTemplate créé : {OUTPUT_PPT}")
    print(f"Tables traités : {table_count}")
    print(f"En-têtes : {NEW_MONTHS[0]} → {NEW_MONTHS[-1]}")


if __name__ == "__main__":
    main()
