#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
tools/ppt_generator.py
-----------------------
Génère le reporting WebPerf PPT à partir des données Google Sheets.
Lit le template, remplit les tableaux de données, sauvegarde le résultat.

Usage:
    python tools/ppt_generator.py [YYYY-MM]     # mois cible (pour le nom du fichier)
    python tools/ppt_generator.py               # = mois courant
    python tools/ppt_generator.py 2026-02
"""

import sys
import os
import datetime
import unicodedata
import gspread
from google.oauth2.service_account import Credentials
from pptx import Presentation
from pptx.oxml.ns import qn


# ── Configuration ──────────────────────────────────────────────────────────────
TEMPLATE_PPT     = "tools/webperf_template.pptx"
CREDENTIALS_FILE = "tools/credentials/service_account.json"
SPREADSHEET_ID   = "1D7IYLK2GQ77L8o-mXJzFzFqxgM0m7DLH0cxtw8khtn8"
OUTPUT_DIR       = "output/webperf"

SCOPES = [
    "https://www.googleapis.com/auth/spreadsheets",
    "https://www.googleapis.com/auth/drive",
]

# ── Mapping slides PPT → feuille Sheets + bloc DSK/MOB ────────────────────────
# slide_idx (1-based) → (sheet_name, "DSK" ou "MOB")
SLIDE_DATA_MAP = {
    13: ("HP",                "DSK"),
    15: ("HP",                "MOB"),
    22: ("HP croisierenet",   "DSK"),
    24: ("HP croisierenet",   "MOB"),
    26: ("HP croisieres.fr",  "DSK"),
    28: ("HP croisieres.fr",  "MOB"),
    30: ("HP croisieres.com", "DSK"),
    32: ("HP croisieres.com", "MOB"),
    39: ("MSC",               "DSK"),
    41: ("MSC",               "MOB"),
    43: ("COSTA",             "DSK"),
    45: ("COSTA",             "MOB"),
    48: ("SL",                "DSK"),
    50: ("SL",                "MOB"),
    53: ("FP",                "DSK"),
    55: ("FP",                "MOB"),
    58: ("LP navire",         "DSK"),
    60: ("LP navire",         "MOB"),
}

# HP concurrent : traitement séparé (scores uniquement, 5r x 15c)
HP_CONCURRENT_SLIDES = {18: "DSK", 20: "MOB"}

# Mapping : ligne PPT table (1-based) → clé métrique dans le bloc Sheets
# (row 0 = headers, rows 1-8 = métriques)
PPT_ROW_TO_METRIC = {
    1: "score",
    2: "fcp",
    3: "lcp",
    4: "cls",
    # row 5 = FID (déprécié) → skip
    6: "si",
    # row 7 = TTI (déprécié) → skip
    8: "tbt",
}

# Layout des lignes dans le Google Sheets (1-indexé)
DSK_HEADER_ROW = 3   # row contenant les mois (0-indexed dans all_values = 2)
DSK_METRIC_ROWS = {  # métrique → index dans all_values (0-based)
    "score": 3,   # Sheets row 4
    "fcp":   4,
    "lcp":   5,
    "cls":   6,
    # row 7 (idx 7) = FID → skip
    "si":    8,   # Sheets row 9
    # row 9 (idx 9) = TTI → skip
    "tbt":   10,  # Sheets row 11
}

MOB_HEADER_ROW = 13
MOB_METRIC_ROWS = {
    "score": 13,  # Sheets row 14
    "fcp":   14,
    "lcp":   15,
    "cls":   16,
    "si":    18,  # Sheets row 19
    "tbt":   20,  # Sheets row 21
}

# HP concurrent (0-based dans all_values)
HP_CONC_DSK_HEADER_IDX = 3   # row 4 dans Sheets
HP_CONC_DSK_DATA_START  = 4  # rows 5-8
HP_CONC_MOB_HEADER_IDX  = 9  # row 10
HP_CONC_MOB_DATA_START  = 10 # rows 11-14


# ── Helpers ────────────────────────────────────────────────────────────────────

def _strip_accents(text):
    return ''.join(
        c for c in unicodedata.normalize('NFD', text)
        if unicodedata.category(c) != 'Mn'
    )


MONTH_NAMES = {
    "janv": 1, "jan": 1, "janvier": 1, "january": 1,
    "fevr": 2, "fev": 2, "fevrier": 2, "february": 2, "feb": 2,
    "mars": 3, "march": 3, "mar": 3,
    "avr": 4, "avril": 4, "april": 4, "apr": 4,
    "mai": 5, "may": 5,
    "juin": 6, "june": 6, "jun": 6,
    "juil": 7, "juillet": 7, "july": 7, "jul": 7,
    "aout": 8, "august": 8, "aug": 8,
    "sept": 9, "sep": 9, "septembre": 9, "september": 9,
    "oct": 10, "octobre": 10, "october": 10,
    "nov": 11, "novembre": 11, "november": 11,
    "dec": 12, "decembre": 12, "december": 12,
}


def parse_month_key(cell_str):
    """
    Parse une cellule contenant un mois et retourne (year, month) ou None.
    Gère : "janv.-26", "févr.-26", "janv-25", "fév-25", etc.
    """
    if not cell_str:
        return None
    clean = _strip_accents(cell_str.lower()).replace('.', '').replace(',', '').strip()
    # Format "mois-annee" : "janv-26", "fevr-26"
    if '-' in clean:
        parts = clean.split('-')
        if len(parts) == 2:
            m_str, y_str = parts[0].strip(), parts[1].strip()
            m_num = MONTH_NAMES.get(m_str)
            try:
                y_raw = int(y_str)
                y_num = 2000 + y_raw if y_raw < 100 else y_raw
                if m_num:
                    return (y_num, m_num)
            except ValueError:
                pass
    # Format "mois annee" : "janv 2026"
    parts = clean.split()
    if len(parts) == 2:
        for m_str, y_str in [(parts[0], parts[1]), (parts[1], parts[0])]:
            m_num = MONTH_NAMES.get(m_str)
            try:
                y_num = int(y_str)
                if m_num:
                    return (y_num, m_num)
            except ValueError:
                continue
    return None


def format_value(val):
    """
    Formate une valeur numérique pour affichage en PPT.
    Utilise la virgule comme séparateur décimal (format français).
    """
    if val is None or val == '':
        return ''
    s = str(val).strip()
    if not s:
        return ''
    # Essayer de normaliser en float
    try:
        f = float(s.replace(',', '.'))
        # Entier → pas de décimale
        if f == int(f) and '.' not in s and ',' not in s:
            return str(int(f))
        # Float → virgule comme séparateur
        return str(f).replace('.', ',')
    except (ValueError, OverflowError):
        return s


def set_cell_text(cell, new_text):
    """Remplace le texte d'une cellule en préservant la mise en forme."""
    t_elems = cell.text_frame._txBody.findall('.//' + qn('a:t'))
    if not t_elems:
        cell.text = new_text
        return
    t_elems[0].text = new_text
    for t_elem in t_elems[1:]:
        t_elem.text = ''


def get_table(slide):
    """Retourne la première table trouvée dans un slide, ou None."""
    for shape in slide.shapes:
        if shape.shape_type == 19:
            return shape.table
    return None


# ── Lecture des données Sheets ─────────────────────────────────────────────────

def read_sheet_data(ws):
    """
    Lit toutes les valeurs d'une feuille et retourne un dict :
    {
      "DSK": {"months": [(year,month), ...], "score": [...], "fcp": [...], ...},
      "MOB": { ... }
    }
    """
    values = ws.get_all_values()

    def get_row(idx):
        return values[idx] if idx < len(values) else []

    def parse_months(header_row):
        months = []
        for cell in header_row:
            key = parse_month_key(cell)
            months.append(key)  # None si pas un mois
        return months

    def get_metric_values(row_idx):
        row = get_row(row_idx)
        return row

    # DSK
    dsk_months = parse_months(get_row(DSK_HEADER_ROW - 1))  # 0-based
    dsk_data = {"months": dsk_months}
    for metric, row_idx in DSK_METRIC_ROWS.items():
        dsk_data[metric] = get_row(row_idx)

    # MOB
    mob_months = parse_months(get_row(MOB_HEADER_ROW - 1))
    mob_data = {"months": mob_months}
    for metric, row_idx in MOB_METRIC_ROWS.items():
        mob_data[metric] = get_row(row_idx)

    return {"DSK": dsk_data, "MOB": mob_data}


def read_concurrent_data(ws):
    """
    Lit la feuille HP concurrent.
    Retourne :
    {
      "DSK": {"months": [...], "rows": [scores_site1, scores_site2, ...]},
      "MOB": { ... }
    }
    """
    values = ws.get_all_values()

    def get_row(idx):
        return values[idx] if idx < len(values) else []

    def parse_months(header_row):
        return [parse_month_key(c) for c in header_row]

    dsk_header = get_row(HP_CONC_DSK_HEADER_IDX)
    mob_header = get_row(HP_CONC_MOB_HEADER_IDX)

    dsk_rows = [get_row(HP_CONC_DSK_DATA_START + i) for i in range(4)]
    mob_rows = [get_row(HP_CONC_MOB_DATA_START + i) for i in range(4)]

    return {
        "DSK": {"months": parse_months(dsk_header), "rows": dsk_rows},
        "MOB": {"months": parse_months(mob_header), "rows": mob_rows},
    }


# ── Remplissage des tableaux PPT ───────────────────────────────────────────────

def fill_metrics_table(table, sheet_data, block):
    """
    Remplit un tableau de métriques 9r x 16c.
    table     : objet pptx Table
    sheet_data: dict retourné par read_sheet_data()
    block     : "DSK" ou "MOB"
    """
    data = sheet_data[block]
    sheet_months = data["months"]  # liste de (year, month) ou None

    # Row 0 = headers de mois dans le PPT
    ppt_header = [parse_month_key(cell.text.strip()) for cell in table.rows[0].cells]

    for col_idx, ppt_month_key in enumerate(ppt_header):
        if ppt_month_key is None:
            continue  # colonne label (Reco, DSK/MOB) → skip

        # Trouver la colonne correspondante dans le Sheets
        sheet_col = None
        for s_idx, s_month in enumerate(sheet_months):
            if s_month == ppt_month_key:
                sheet_col = s_idx
                break

        # Remplir chaque ligne de métrique
        for ppt_row_idx, metric_key in PPT_ROW_TO_METRIC.items():
            if ppt_row_idx >= len(table.rows):
                continue
            cell = table.rows[ppt_row_idx].cells[col_idx]

            if sheet_col is None:
                set_cell_text(cell, '')
                continue

            row_data = data.get(metric_key, [])
            val = row_data[sheet_col] if sheet_col < len(row_data) else ''
            set_cell_text(cell, format_value(val))


def fill_concurrent_table(table, conc_data, block):
    """
    Remplit un tableau HP concurrent 5r x 15c.
    """
    data = conc_data[block]
    sheet_months = data["months"]
    site_rows = data["rows"]  # liste de 4 listes de valeurs

    ppt_header = [parse_month_key(cell.text.strip()) for cell in table.rows[0].cells]

    for col_idx, ppt_month_key in enumerate(ppt_header):
        if ppt_month_key is None:
            continue

        sheet_col = None
        for s_idx, s_month in enumerate(sheet_months):
            if s_month == ppt_month_key:
                sheet_col = s_idx
                break

        # Remplir les 4 lignes de sites (rows 1-4)
        for site_idx in range(4):
            ppt_row_idx = site_idx + 1
            if ppt_row_idx >= len(table.rows):
                continue
            cell = table.rows[ppt_row_idx].cells[col_idx]

            if sheet_col is None or site_idx >= len(site_rows):
                set_cell_text(cell, '')
                continue

            row_data = site_rows[site_idx]
            val = row_data[sheet_col] if sheet_col < len(row_data) else ''
            set_cell_text(cell, format_value(val))


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) > 1:
        year_month = sys.argv[1]
    else:
        year_month = datetime.datetime.now().strftime("%Y-%m")

    try:
        target_dt = datetime.datetime.strptime(year_month, "%Y-%m")
    except ValueError:
        print(f"Format invalide : '{year_month}'. Utiliser YYYY-MM")
        sys.exit(1)

    output_file = os.path.join(OUTPUT_DIR, f"webperf_AB_{target_dt.strftime('%m-%Y')}.pptx")
    print(f"PPT Generator — {target_dt.strftime('%B %Y')}")
    print(f"Template : {TEMPLATE_PPT}")
    print(f"Output   : {output_file}")
    print()

    # ── Connexion Google Sheets ───────────────────────────────────────────────
    print("Connexion Google Sheets...")
    creds = Credentials.from_service_account_file(CREDENTIALS_FILE, scopes=SCOPES)
    gc    = gspread.authorize(creds)
    sh    = gc.open_by_key(SPREADSHEET_ID)
    print(f"Connecté : '{sh.title}'\n")

    # ── Lecture de toutes les feuilles ────────────────────────────────────────
    print("Lecture des données Sheets...")
    sheet_names = set(v[0] for v in SLIDE_DATA_MAP.values())
    sheets_data = {}
    for name in sheet_names:
        try:
            ws = sh.worksheet(name)
            sheets_data[name] = read_sheet_data(ws)
            print(f"  {name} — OK")
        except Exception as e:
            print(f"  {name} — ERREUR : {e}")

    # HP concurrent
    try:
        ws_conc = sh.worksheet("HP concurrent")
        conc_data = read_concurrent_data(ws_conc)
        print(f"  HP concurrent — OK")
    except Exception as e:
        print(f"  HP concurrent — ERREUR : {e}")
        conc_data = None
    print()

    # ── Ouverture du template ─────────────────────────────────────────────────
    print(f"Ouverture du template : {TEMPLATE_PPT}")
    prs = Presentation(TEMPLATE_PPT)

    # ── Remplissage des tableaux ──────────────────────────────────────────────
    print("Remplissage des tableaux...")
    filled = 0

    for slide_idx_1based, (sheet_name, block) in SLIDE_DATA_MAP.items():
        slide = prs.slides[slide_idx_1based - 1]
        table = get_table(slide)
        if table is None:
            print(f"  Slide {slide_idx_1based:2d} — pas de table")
            continue
        if sheet_name not in sheets_data:
            print(f"  Slide {slide_idx_1based:2d} — données manquantes pour '{sheet_name}'")
            continue

        fill_metrics_table(table, sheets_data[sheet_name], block)
        print(f"  Slide {slide_idx_1based:2d} — {sheet_name} {block} OK")
        filled += 1

    # HP concurrent
    if conc_data:
        for slide_idx_1based, block in HP_CONCURRENT_SLIDES.items():
            slide = prs.slides[slide_idx_1based - 1]
            table = get_table(slide)
            if table:
                fill_concurrent_table(table, conc_data, block)
                print(f"  Slide {slide_idx_1based:2d} — HP concurrent {block} OK")
                filled += 1

    print()

    # ── Sauvegarde ────────────────────────────────────────────────────────────
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    prs.save(output_file)
    print(f"PPT généré : {output_file}")
    print(f"Slides remplis : {filled}/{len(SLIDE_DATA_MAP) + len(HP_CONCURRENT_SLIDES)}")


if __name__ == "__main__":
    main()
